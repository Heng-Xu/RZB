from __future__ import annotations

import csv
import json
import math
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from typing import NamedTuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "reports" / "06.10"
OUT = OUT_DIR / "battery_trajectory_group_meeting_20260610.pptx"
SUMMARY_OUT = OUT_DIR / "ppt_generation_summary.json"

DATASETS = ("HUST", "XJTU", "MIT", "TJU")
WINDOWS = ("010", "020", "040", "080")
WINDOW_LABEL = "10/20/40/80"
STAGE_BY_WINDOW = {win: f"fixed{win}" for win in WINDOWS}

FULL_MODEL_SPECS = (
    ("Hybrid 主线", "Four_models", "mainline_rebuild_unified"),
    ("LSTM", "Four_models", "lstm_full_trajectory"),
    ("LSTM-PINN", "Four_models", "lstm_pinn_ft_full_trajectory"),
    ("xPatch-DT", "Four_models", "xpatch_dt_fixed_history"),
    ("BatteryGPT", "BatteryGPT", "batterygpt_xpatchflow_scratch"),
)

METRICS = (
    ("rmse_mean", "RMSE\nmean", "global_rmse_mean", 4, "min", True),
    ("rmse_std", "RMSE\nstd", "global_rmse_std", 4, "min", False),
    ("mape_mean", "MAPE\nmean", "global_mape_mean", 2, "min", True),
    ("mape_std", "MAPE\nstd", "global_mape_std", 2, "min", False),
)

COLORS = {
    "ink": RGBColor(27, 35, 45),
    "muted": RGBColor(91, 103, 115),
    "navy": RGBColor(19, 41, 73),
    "blue": RGBColor(35, 95, 166),
    "green": RGBColor(42, 126, 87),
    "red": RGBColor(175, 58, 62),
    "orange": RGBColor(205, 123, 47),
    "light": RGBColor(247, 249, 252),
    "line": RGBColor(204, 214, 226),
    "pale_blue": RGBColor(232, 241, 252),
    "pale_green": RGBColor(232, 246, 239),
    "pale_orange": RGBColor(254, 240, 222),
    "white": RGBColor(255, 255, 255),
}


class MetricPart(NamedTuple):
    text: str
    best: bool = False


class MetricCell(NamedTuple):
    parts: tuple[MetricPart, ...]


class SummaryRow(NamedTuple):
    label: str
    metric_cells: dict[str, MetricCell]


class ResultTable(NamedTuple):
    dataset: str
    rows: list[SummaryRow]
    source_files: dict[str, dict[str, str]]
    seed_counts: dict[str, dict[str, int]]


def log_runtime() -> dict:
    probe = (
        "import json, subprocess, sys\n"
        "info={'python': sys.executable, 'torch_cuda_available': None, "
        "'torch_cuda_device_count': None, 'nvidia_smi_gpus': [], "
        "'probe_mode': 'isolated_python'}\n"
        "try:\n"
        "    smi=subprocess.run(['nvidia-smi','--query-gpu=index,name,memory.used,memory.total',"
        "'--format=csv,noheader'], check=True, capture_output=True, text=True)\n"
        "    info['nvidia_smi_gpus']=[line.strip() for line in smi.stdout.splitlines() if line.strip()]\n"
        "except Exception as exc:\n"
        "    info['nvidia_smi_error']=repr(exc)\n"
        "try:\n"
        "    import torch\n"
        "    info['torch_cuda_available']=bool(torch.cuda.is_available())\n"
        "    info['torch_cuda_device_count']=int(torch.cuda.device_count())\n"
        "except Exception as exc:\n"
        "    info['torch_cuda_error']=repr(exc)\n"
        "print(json.dumps(info, ensure_ascii=False))\n"
    )
    try:
        result = subprocess.run(
            [sys.executable, "-c", probe],
            check=True,
            capture_output=True,
            text=True,
        )
        info = json.loads(result.stdout)
    except Exception as exc:
        info = {
            "python": sys.executable,
            "torch_cuda_available": None,
            "torch_cuda_device_count": None,
            "nvidia_smi_gpus": [],
            "probe_mode": "isolated_python_failed",
            "probe_error": repr(exc),
        }
    print(f"python={info['python']}")
    print(f"torch_cuda_available={info.get('torch_cuda_available')}")
    print(f"torch_cuda_device_count={info.get('torch_cuda_device_count')}")
    print(f"nvidia_smi_gpus={info.get('nvidia_smi_gpus')}")
    return info


def read_compare_csv(path: Path) -> list[dict[str, str]]:
    if not path.exists():
        raise FileNotFoundError(path)
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def candidate_summary_files(dataset: str, family: str) -> list[Path]:
    result_dir = ROOT / "compare" / "results" / dataset
    return sorted(result_dir.glob(f"{family}_*_seed_summary.csv"), reverse=True)


def candidate_ablation_files(dataset: str) -> list[Path]:
    result_dir = ROOT / "compare" / "results" / "BV_ablation" / dataset
    return sorted(result_dir.glob("BV_ablation_*_seed_summary.csv"), reverse=True)


def find_summary_row(
    *,
    dataset: str,
    family: str,
    window: str,
    model: str,
    ablation: bool = False,
) -> tuple[dict[str, str], Path]:
    stage = STAGE_BY_WINDOW[window]
    files = candidate_ablation_files(dataset) if ablation else candidate_summary_files(dataset, family)
    for path in files:
        for row in read_compare_csv(path):
            if row.get("stage_tag") == stage and row.get("model") == model:
                return row, path
    raise FileNotFoundError(f"No {dataset} {family} {stage} {model} summary row")


def row_float(row: dict[str, str], key: str) -> float | None:
    value = row.get(key)
    try:
        parsed = float(value) if value not in (None, "") else None
    except ValueError:
        return None
    if parsed is None or math.isnan(parsed):
        return None
    return parsed


def row_int(row: dict[str, str], key: str) -> int:
    value = row.get(key)
    try:
        return int(float(value)) if value not in (None, "") else 0
    except ValueError:
        return 0


def fmt_metric(value: float | None, digits: int) -> str:
    if value is None:
        return "-"
    return f"{value:.{digits}f}"


def mark_best(values_by_label: dict[str, list[float | None]], *, direction: str) -> dict[str, list[bool]]:
    flags = {label: [False] * len(WINDOWS) for label in values_by_label}
    for idx, _win in enumerate(WINDOWS):
        candidates = {
            label: values[idx]
            for label, values in values_by_label.items()
            if values[idx] is not None
        }
        if not candidates:
            continue
        best = min(candidates.values()) if direction == "min" else max(candidates.values())
        for label, value in candidates.items():
            flags[label][idx] = math.isclose(value, best, rel_tol=0.0, abs_tol=1e-12)
    return flags


def build_metric_cells(raw_rows: dict[str, list[dict[str, str]]]) -> dict[str, dict[str, MetricCell]]:
    metric_cells = {label: {} for label in raw_rows}
    for metric_key, _header, csv_key, digits, direction, mark_best_value in METRICS:
        values = {
            label: [row_float(row, csv_key) for row in rows]
            for label, rows in raw_rows.items()
        }
        flags = mark_best(values, direction=direction) if mark_best_value else {
            label: [False] * len(WINDOWS) for label in raw_rows
        }
        for label, vals in values.items():
            metric_cells[label][metric_key] = MetricCell(
                tuple(
                    MetricPart(fmt_metric(value, digits), flags[label][idx])
                    for idx, value in enumerate(vals)
                )
            )
    return metric_cells


def build_full_table(dataset: str) -> ResultTable:
    raw_rows: dict[str, list[dict[str, str]]] = {}
    sources: dict[str, dict[str, str]] = {}
    seed_counts: dict[str, dict[str, int]] = {}
    for label, family, model in FULL_MODEL_SPECS:
        raw_rows[label] = []
        sources[label] = {}
        seed_counts[label] = {}
        for win in WINDOWS:
            row, path = find_summary_row(dataset=dataset, family=family, window=win, model=model)
            raw_rows[label].append(row)
            sources[label][win] = str(path.relative_to(ROOT))
            seed_counts[label][win] = row_int(row, "global_rmse_count")
    metric_cells = build_metric_cells(raw_rows)
    rows = [
        SummaryRow(label=label, metric_cells=metric_cells[label])
        for label, _family, _model in FULL_MODEL_SPECS
    ]
    return ResultTable(dataset=dataset, rows=rows, source_files=sources, seed_counts=seed_counts)


def build_ablation_table(dataset: str) -> ResultTable:
    raw_rows = {"主线": [], "去 BV 损失": []}
    sources = {"主线": {}, "去 BV 损失": {}}
    seed_counts = {"主线": {}, "去 BV 损失": {}}
    for win in WINDOWS:
        mainline_row, mainline_path = find_summary_row(
            dataset=dataset,
            family="Four_models",
            window=win,
            model="mainline_rebuild_unified",
        )
        ablation_row, ablation_path = find_summary_row(
            dataset=dataset,
            family="BV_ablation",
            window=win,
            model="mainline_rebuild_unified",
            ablation=True,
        )
        raw_rows["主线"].append(mainline_row)
        raw_rows["去 BV 损失"].append(ablation_row)
        sources["主线"][win] = str(mainline_path.relative_to(ROOT))
        sources["去 BV 损失"][win] = str(ablation_path.relative_to(ROOT))
        seed_counts["主线"][win] = row_int(mainline_row, "global_rmse_count")
        seed_counts["去 BV 损失"][win] = row_int(ablation_row, "global_rmse_count")
    metric_cells = build_metric_cells(raw_rows)
    rows = [
        SummaryRow(label=label, metric_cells=metric_cells[label])
        for label in ("主线", "去 BV 损失")
    ]
    return ResultTable(dataset=dataset, rows=rows, source_files=sources, seed_counts=seed_counts)


def set_run_font(
    run,
    *,
    size: int,
    bold: bool = False,
    color: RGBColor | None = None,
    name: str = "Microsoft YaHei",
) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 10,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN | None = None,
) -> object:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    *,
    line: RGBColor | None = None,
    radius: MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
) -> object:
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.color.rgb = fill
        shape.line.transparency = 100000
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_rule(slide, x: float, y: float, w: float, *, color: RGBColor = COLORS["line"]):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = color
    line.line.width = Pt(0.8)
    return line


def add_title(slide, title: str, subtitle: str = "") -> None:
    add_textbox(slide, 0.42, 0.20, 9.4, 0.42, title, size=20, bold=True, color=COLORS["navy"])
    if subtitle:
        add_textbox(slide, 0.45, 0.62, 10.8, 0.30, subtitle, size=10, color=COLORS["muted"])
    add_textbox(slide, 12.15, 0.28, 0.75, 0.22, "06.10", size=8, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_bullets(slide, x: float, y: float, w: float, h: float, lines: list[str], *, size: int = 10):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.space_after = Pt(4)
        p.line_spacing = 1.05
        for run in p.runs:
            set_run_font(run, size=size, color=COLORS["ink"])
    return box


def clear_cell(cell, *, fill: RGBColor | None = None, valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE):
    cell.text = ""
    cell.margin_left = Inches(0.025)
    cell.margin_right = Inches(0.025)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    cell.vertical_anchor = valign
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return tf.paragraphs[0]


def cell_plain(
    cell,
    text: str,
    *,
    size: int = 8,
    bold: bool = False,
    color: RGBColor = COLORS["ink"],
    align: PP_ALIGN = PP_ALIGN.CENTER,
    fill: RGBColor | None = None,
):
    p = clear_cell(cell, fill=fill)
    p.alignment = align
    p.line_spacing = 0.88
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color)


def cell_metric(cell, metric_cell: MetricCell, *, fill: RGBColor | None = None, size: int = 7):
    p = clear_cell(cell, fill=fill)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 0.86
    for idx, part in enumerate(metric_cell.parts):
        if idx:
            sep = p.add_run()
            sep.text = "/"
            set_run_font(sep, size=size, color=COLORS["muted"])
        run = p.add_run()
        run.text = part.text
        set_run_font(run, size=size, bold=part.best, color=COLORS["blue"] if part.best else COLORS["ink"])


def add_compact_table_shape(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    rows_count: int,
    *,
    accent: RGBColor,
):
    cols = 5
    table_shape = slide.shapes.add_table(rows_count, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    widths = [1.34, 1.13, 1.02, 1.13, 1.02]
    scale = w / sum(widths)
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width * scale)
    table.rows[0].height = Inches(0.32)
    body_h = max(0.23, (h - 0.32) / (rows_count - 1))
    for idx in range(1, rows_count):
        table.rows[idx].height = Inches(body_h)
    headers = ["模型", "RMSE\nmean", "RMSE\nstd", "MAPE\nmean", "MAPE\nstd"]
    for col, header in enumerate(headers):
        cell_plain(table.cell(0, col), header, size=7, bold=True, color=COLORS["white"], fill=accent)
    return table_shape


def add_result_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    table_data: ResultTable,
    *,
    title: str,
    accent: RGBColor,
):
    add_textbox(slide, x, y, w, 0.24, title, size=12, bold=True, color=COLORS["navy"])
    add_textbox(slide, x, y + 0.27, w, 0.18, f"窗口顺序 {WINDOW_LABEL}；5 seeds/splits", size=7, color=COLORS["muted"])
    table_shape = add_compact_table_shape(
        slide,
        x,
        y + 0.50,
        w,
        h - 0.50,
        len(table_data.rows) + 1,
        accent=accent,
    )
    table = table_shape.table
    for row_idx, row in enumerate(table_data.rows, start=1):
        fill = COLORS["white"] if row_idx % 2 else COLORS["light"]
        cell_plain(table.cell(row_idx, 0), row.label, size=7, bold=True, align=PP_ALIGN.LEFT, fill=fill)
        for col_idx, (metric_key, _header, _csv_key, _digits, _direction, _mark) in enumerate(METRICS, start=1):
            cell_metric(table.cell(row_idx, col_idx), row.metric_cells[metric_key], fill=fill, size=7)
    return table_shape


def slide_full_results(prs: Presentation, tables: dict[str, ResultTable]) -> None:
    slide = add_blank_slide(prs)
    add_title(
        slide,
        "数据划分后全量结果",
        "固定窗口 10/20/40/80；表格内按窗口顺序展示 RMSE/MAPE 的均值与标准差，均值最优加粗",
    )
    positions = {
        "HUST": (0.45, 1.03, 6.02, 2.72, COLORS["green"]),
        "XJTU": (6.86, 1.03, 6.02, 2.72, COLORS["blue"]),
        "MIT": (0.45, 4.12, 6.02, 2.72, COLORS["navy"]),
        "TJU": (6.86, 4.12, 6.02, 2.72, COLORS["orange"]),
    }
    for dataset, (x, y, w, h, accent) in positions.items():
        add_result_table(slide, x, y, w, h, tables[dataset], title=f"{dataset} 全量对比结果", accent=accent)


def slide_ablation(prs: Presentation, tables: dict[str, ResultTable]) -> None:
    slide = add_blank_slide(prs)
    add_title(
        slide,
        "物理损失消融结果",
        "主线模型与去 BV 损失对照；表格内按窗口顺序展示 RMSE/MAPE 的均值与标准差，均值更优加粗",
    )
    positions = {
        "HUST": (0.45, 1.03, 6.02, 2.30, COLORS["green"]),
        "XJTU": (6.86, 1.03, 6.02, 2.30, COLORS["blue"]),
        "MIT": (0.45, 4.02, 6.02, 2.30, COLORS["navy"]),
        "TJU": (6.86, 4.02, 6.02, 2.30, COLORS["orange"]),
    }
    for dataset, (x, y, w, h, accent) in positions.items():
        add_result_table(slide, x, y, w, h, tables[dataset], title=f"{dataset} 主线 BV 消融", accent=accent)
    add_textbox(
        slide,
        0.55,
        6.70,
        12.20,
        0.22,
        "注：去 BV 损失仅作为物理损失消融对照，不改变数据划分、窗口与指标统计口径。",
        size=8,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )


PAPER_SUMMARY = (
    "本项目围绕公开电池退化数据集的早期轨迹预测，构建融合仿真数据、域适应、BV 物理损失与不确定性估计的 Hybrid 主线模型；"
    "在 HUST、XJTU、MIT、TJU 四个数据集上统一电池级划分、固定历史窗口和跨 seed 统计口径，并通过 BatteryGPT、LSTM、LSTM-PINN、xPatch-DT 等对比方法验证模型精度与物理约束贡献。"
)

PAPER_OUTLINE = [
    "1 引言：电池健康轨迹预测需求、公开数据集泛化挑战、物理约束学习的研究动机。",
    "2 相关工作：数据驱动 SOH/RUL 预测、物理信息神经网络、域适应与大模型序列建模方法。",
    "3 方法：Hybrid 主干、仿真-真实数据融合、域适应对齐、BV 残差损失、不确定性估计。",
    "4 实验设置：HUST/XJTU/MIT/TJU 数据集，电池级划分，固定历史窗口，baseline 与消融方案。",
    "5 结果与分析：全量对比、BV 损失消融、跨窗口稳定性、误差来源与典型轨迹案例。",
    "6 结论：总结方法收益，讨论数据集迁移限制与后续多工况扩展。",
]


def slide_paper_outline(prs: Presentation) -> None:
    slide = add_blank_slide(prs)
    add_title(slide, "论文大纲写作", "左侧为初稿文字；右侧留白用于插入总体框架图或结果示意图")

    add_rect(slide, 0.55, 1.08, 6.00, 5.90, COLORS["white"], line=COLORS["line"], radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_textbox(slide, 0.78, 1.28, 5.52, 0.28, "项目总结", size=14, bold=True, color=COLORS["navy"])
    add_rule(slide, 0.78, 1.68, 5.52)
    add_textbox(slide, 0.80, 1.88, 5.48, 1.05, PAPER_SUMMARY, size=10, color=COLORS["ink"])
    add_textbox(slide, 0.78, 3.18, 5.52, 0.28, "论文大纲初稿", size=14, bold=True, color=COLORS["navy"])
    add_rule(slide, 0.78, 3.58, 5.52)
    add_bullets(slide, 0.80, 3.78, 5.44, 2.72, PAPER_OUTLINE, size=8)

    add_rect(slide, 6.98, 1.08, 5.80, 5.90, COLORS["light"], line=COLORS["line"], radius=MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    add_textbox(slide, 7.30, 3.58, 5.18, 0.26, "图片预留区", size=16, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.52, 4.00, 4.75, 0.22, "可放置方法框架图 / 四数据集结果示意图", size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def table_to_summary(table: ResultTable) -> dict:
    return {
        "dataset": table.dataset,
        "rows": [
            {
                "label": row.label,
                "metrics": {
                    key: [{"text": part.text, "best": part.best} for part in cell.parts]
                    for key, cell in row.metric_cells.items()
                },
            }
            for row in table.rows
        ],
        "source_files": table.source_files,
        "seed_counts": table.seed_counts,
    }


def build_ppt() -> dict:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    runtime = log_runtime()
    full_tables = {dataset: build_full_table(dataset) for dataset in DATASETS}
    ablation_tables = {dataset: build_ablation_table(dataset) for dataset in DATASETS}

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    slide_full_results(prs, full_tables)
    slide_ablation(prs, ablation_tables)
    slide_paper_outline(prs)

    prs.save(OUT)

    summary = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "output": str(OUT.relative_to(ROOT)),
        "runtime": runtime,
        "slides": 3,
        "font_levels": [20, 14, 12, 10, 8, 7],
        "metric_windows": WINDOW_LABEL,
        "metric_columns": [metric_key for metric_key, *_rest in METRICS],
        "full_tables": {dataset: table_to_summary(table) for dataset, table in full_tables.items()},
        "ablation_tables": {dataset: table_to_summary(table) for dataset, table in ablation_tables.items()},
        "paper_outline": {
            "summary_text": PAPER_SUMMARY,
            "outline": PAPER_OUTLINE,
            "has_image_placeholder": True,
        },
    }
    SUMMARY_OUT.write_text(json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"wrote={OUT}")
    print(f"summary={SUMMARY_OUT}")
    return summary


if __name__ == "__main__":
    build_ppt()
