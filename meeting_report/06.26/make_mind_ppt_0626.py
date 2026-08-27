#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""组会汇报 PPT（3 页，简明）：FG-NSGA-II —— 县级容载比弹性的可行性引导多目标进化优化。
重点＝**优化算法上的核心创新**，用对比图 + 对比表证明；应用结果压成 1 页。
数据口径：收敛扫描 run beyb8rsgp（K∈{10,20,30,40}，5 变体逐代 HV），主算例 K=20。
风格沿用 06.10/06.26 版助手（同配色/字体/版式）。"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "report" / "06.26"
OUT = OUT_DIR / "容载比EA_组会汇报_20260626.pptx"
FIG = ROOT / "实验" / "有EA" / "results" / "figures"

COLORS = {
    "ink": RGBColor(27, 35, 45), "muted": RGBColor(91, 103, 115),
    "navy": RGBColor(19, 41, 73), "blue": RGBColor(35, 95, 166),
    "green": RGBColor(42, 126, 87), "red": RGBColor(175, 58, 62),
    "orange": RGBColor(205, 123, 47), "light": RGBColor(247, 249, 252),
    "line": RGBColor(204, 214, 226), "pale_blue": RGBColor(232, 241, 252),
    "pale_green": RGBColor(232, 246, 239), "pale_orange": RGBColor(254, 240, 222),
    "pale_red": RGBColor(250, 234, 234), "white": RGBColor(255, 255, 255),
}


def set_run_font(run, *, size, bold=False, color=None, name="Microsoft YaHei"):
    run.font.name = name; run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, x, y, w, h, text, *, size=10, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run(); run.text = text
    set_run_font(run, size=size, bold=bold, color=color)
    return box


def add_rect(slide, x, y, w, h, fill, *, line=None,
             radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.color.rgb = fill; shape.line.transparency = 100000
    else:
        shape.line.color.rgb = line; shape.line.width = Pt(1.0)
    shape.shadow.inherit = False
    return shape


def add_title(slide, title, subtitle=""):
    add_textbox(slide, 0.45, 0.22, 12.4, 0.5, title, size=22, bold=True, color=COLORS["navy"])
    if subtitle:
        add_textbox(slide, 0.47, 0.74, 12.4, 0.32, subtitle, size=11, color=COLORS["muted"])
    add_textbox(slide, 12.05, 0.30, 0.9, 0.22, "MIND 2026", size=8,
                color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, lines, *, size=11, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sub = line.startswith("  ")
        p.text = ("◦ " if sub else "• ") + line.strip()
        p.space_after = Pt(3); p.line_spacing = 1.08
        for run in p.runs:
            set_run_font(run, size=size - (1 if sub else 0),
                         color=color or (COLORS["muted"] if sub else COLORS["ink"]))
    return box


def box_header(slide, x, y, w, text, fill, tcolor):
    add_rect(slide, x, y, w, 0.42, fill)
    add_textbox(slide, x + 0.12, y + 0.04, w - 0.24, 0.34, text, size=13, bold=True, color=tcolor)


def add_table(slide, x, y, w, h, data, *, col_w=None, fs=10.5, head_fs=10.5,
              head_fill="navy", zebra=True, first_col_bold=True):
    """原生 PPTX 三色表：首行表头，可选斑马纹；data[0]=表头。"""
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    t = gt.table
    if col_w:
        for ci, cw in enumerate(col_w):
            t.columns[ci].width = Inches(cw)
    for ri, row in enumerate(data):
        t.rows[ri].height = Inches(0.34)
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.margin_left = Inches(0.05); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = COLORS[head_fill]
            elif zebra and ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = COLORS["light"]
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = COLORS["white"]
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            is_head = (ri == 0)
            is_fg = (ci == cols - 3 and not is_head)  # FG (ours) 列强调（见列序）
            set_run_font(r, size=head_fs if is_head else fs, bold=is_head or (first_col_bold and ci == 0),
                         color=COLORS["white"] if is_head else COLORS["ink"])
    return t


# ── Slide 1：方法框架 + 核心创新点 ──────────────────────────────────────────────
def slide_framework(prs):
    s = add_blank_slide(prs)
    add_title(s, "方法框架与核心创新：FG-NSGA-II（可行性引导 NSGA-II）",
              "投稿 MIND 2026（IEEE-CIS / EI）· 纯进化算法 · 公开数据合成县级基准 · 主算例 K=20")
    bw, by, bh = 3.95, 1.35, 4.15
    xs = [0.45, 4.69, 8.93]
    # ① 输入
    add_rect(s, xs[0], by, bw, bh, COLORS["pale_blue"], line=COLORS["blue"])
    box_header(s, xs[0], by, bw, "① 输入", COLORS["blue"], COLORS["white"])
    add_bullets(s, xs[0] + 0.14, by + 0.56, bw - 0.28, bh - 0.66, [
        "每站参数：源荷比、联络度、反送时长",
        "物理评价：双向 Z4-LCC（~4 ms/次，无代理）",
        "县 = 20 座 110kV 站",
        "容载比带 + 上级潮流 → 站间耦合",
    ], size=12.5)
    # ② FG-NSGA-II（标注创新）
    add_rect(s, xs[1], by, bw, bh, COLORS["pale_red"], line=COLORS["red"])
    box_header(s, xs[1], by, bw, "② FG-NSGA-II ★创新", COLORS["red"], COLORS["white"])
    add_bullets(s, xs[1] + 0.14, by + 0.56, bw - 0.28, bh - 0.66, [
        "决策：每站容载比 R + 储能 P",
        "目标：min 年化成本、min N-1 缺供",
        "硬约束：反送闸 · 容载比带 · 上级潮流",
        "★ 修复算子（主）：解拉回可行域，不占预算",
        "★ 热启动（辅）：初始种群即合规",
    ], size=12.5)
    # ③ 输出
    add_rect(s, xs[2], by, bw, bh, COLORS["pale_green"], line=COLORS["green"])
    box_header(s, xs[2], by, bw, "③ 输出", COLORS["green"], COLORS["white"])
    add_bullets(s, xs[2] + 0.14, by + 0.56, bw - 0.28, bh - 0.66, [
        "成本 ↔ N-1 风险 Pareto 前沿",
        "拐点 → 一站一策（每站 R*、P*）",
        "聚合容载比 1.80 < 2.0",
        "相对统一 2.0 省 13.1%",
    ], size=12.5)
    for ax in (xs[0] + bw + 0.03, xs[1] + bw + 0.03):
        add_textbox(s, ax, by + bh / 2 - 0.3, 0.5, 0.6, "➜", size=24, bold=True, color=COLORS["navy"])
    # 底部：核心创新（痛点 → 做法）
    add_rect(s, 0.45, 5.78, 12.43, 1.42, COLORS["light"], line=COLORS["line"])
    add_textbox(s, 0.64, 5.88, 12.1, 0.34, "核心创新：可行性由构造保证，而非惩罚调参", size=14, bold=True, color=COLORS["red"])
    add_bullets(s, 0.62, 6.28, 12.2, 0.9, [
        "痛点：三重硬约束下可行域随规模骤缩——经典 NSGA-II / NSGA-III 前期撞不进可行域，收敛慢。",
        "做法：修复算子把解拉回可行域（主）、热启动让初始种群即合规（辅）——评估预算全用于成本↔风险权衡。",
    ], size=12.5)


# ── Slide 2（核心）：算法创新的对比证据（图 + 表）────────────────────────────────
def slide_algorithm(prs, table_rows):
    s = add_blank_slide(prs)
    add_title(s, "核心创新：同质量下约 5–12× 更快收敛（图 + 表）",
              "5 变体 × K∈{10,20,30,40} × 3 seed 逐代 HV；baseline gen800、FG/消融 gen400；每 K 共同参考点")
    # 左（核心图）：4 个 K 的逐代 HV 收敛对比（5 变体）
    conv = FIG / "fig_enhanced_convergence.png"
    if conv.exists():
        s.shapes.add_picture(str(conv), Inches(0.42), Inches(1.40), width=Inches(5.60))
    add_textbox(s, 0.48, 5.10, 6.0, 0.34,
                "图：逐代 HV 收敛（5 变体）—— FG 与修复最快；终值差距小、已如实披露",
                size=10, color=COLORS["muted"])
    # 右：收敛 HV 对比表
    box_header(s, 6.55, 1.40, 6.33, "表：收敛 HV 比 + 加速比（加速比 = 头条）",
               COLORS["navy"], COLORS["white"])
    header = ["规模 K", "FG /\nClassic", "FG /\nNSGA-III", "Repair-only\n/ FG", "加速×\nClassic", "加速×\nNSGA-III"]
    add_table(s, 6.55, 1.98, 6.33, 1.55, [header] + table_rows,
              col_w=[0.95, 1.08, 1.12, 1.28, 0.95, 0.95], fs=11.5, head_fs=10.5)
    add_textbox(s, 6.57, 3.62, 6.3, 0.4,
                "加速比 = FG 达 baseline gen800 质量所需代数之比（预算无关，全域 4.8–11.8×）。",
                size=10, color=COLORS["muted"])
    add_textbox(s, 6.57, 4.06, 6.3, 0.4,
                "pop=100 固定 → 代数与评价次数、墙钟成正比 → 更少代数 = 更省算力。",
                size=10, color=COLORS["muted"])
    # 底部：三条要点
    add_rect(s, 0.40, 5.50, 12.48, 1.72, COLORS["pale_red"], line=COLORS["red"])
    add_textbox(s, 0.55, 5.58, 8.0, 0.34, "三条结论（诚实口径）", size=14, bold=True, color=COLORS["red"])
    add_bullets(s, 0.58, 6.00, 12.2, 1.15, [
        "效率（头条）：FG 用半预算（gen400）即胜过两 baseline 的 gen800——约 5–12× 更少代数达同等质量。",
        "质量（诚实）：收敛 HV 仅略高（K≤30 <1%、K=40 4.8%），久跑 baseline 追平——主张效率，不吹质量。",
        "消融：修复 = 主驱动（repair-only ≈ FG）；热启动 = 早期加速（前期领先、收敛后被洗掉）。",
    ], size=12.5)


# ── Slide 3：应用结果（简）────────────────────────────────────────────────────
def slide_application(prs, robust_line=""):
    s = add_blank_slide(prs)
    add_title(s, "应用结果：县级“一站一策”弹性方案（K=20，收敛口径）",
              "FG-NSGA-II 拐点方案 · 公开数据合成县级基准 · 相对全县统一容载比 2.0")
    pareto = FIG / "fig_county_pareto.png"
    if pareto.exists():
        s.shapes.add_picture(str(pareto), Inches(0.45), Inches(1.30), height=Inches(3.95))
    strat = FIG / "fig_strategy.png"
    if strat.exists():
        s.shapes.add_picture(str(strat), Inches(6.05), Inches(1.30), height=Inches(3.95))
    add_rect(s, 0.40, 5.50, 12.48, 1.72, COLORS["pale_green"], line=COLORS["green"])
    add_textbox(s, 0.55, 5.58, 8.0, 0.34, "应用结论（已校核）", size=14, bold=True, color=COLORS["green"])
    add_bullets(s, 0.58, 6.00, 6.15, 1.15, [
        "拐点：成本 18882 万/年、N-1 风险 114 MWh、聚合容载比 1.80 < 2.0",
        "一站一策：11 站配储 / 9 站低容载比+强联络 / 0 站增容（最大 R* ≈ 1.96）",
    ], size=12.5)
    add_bullets(s, 6.95, 6.00, 5.85, 1.15, [
        "相对统一容载比 2.0（21836 万）同风险省 13.1%",
        (robust_line or "6 县鲁棒：12.3%±1.3%（10.1–13.7%，全为正）"),
        "验证：小 K 达真前沿；K≥12 不可枚举 → 须用 EA（K=40 单跑 ~4 min）",
    ], size=12.5)


def build():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    # 从收敛扫描 CSV 读取真实对比数据（比值 + 加速比），避免硬编码漂移
    import csv
    conv = ROOT / "实验" / "有EA" / "results" / "ea_converge.csv"
    table_rows = []
    if conv.exists():
        import pandas as pd
        df = pd.read_csv(conv)
        mg = df.groupby("algo").gen.max().to_dict()
        for K in sorted(df.K.unique()):
            s = df[df.K == K]
            def hv(a):
                return s[(s.algo == a) & (s.gen == mg[a])].hv.mean()
            fg, cl, n3, rp = hv("fg"), hv("classic"), hv("nsga3"), hv("repair_only")
            fgc = s[s.algo == "fg"].groupby("gen").hv.mean()
            sp = {}
            for b in ("classic", "nsga3"):
                bt = hv(b); r = fgc[fgc >= bt]
                g = int(r.index.min()) if len(r) else mg["fg"]; sp[b] = mg[b] / g
            table_rows.append([str(int(K)), f"{fg/cl:.3f}", f"{fg/n3:.3f}", f"{rp/fg:.3f}",
                               f"{sp['classic']:.1f}×", f"{sp['nsga3']:.1f}×"])
    # 多县鲁棒性
    robust_line = ""
    rob = ROOT / "实验" / "有EA" / "results" / "ea_robustness.csv"
    if rob.exists():
        rows = list(csv.DictReader(open(rob, encoding="utf-8")))
        if rows:
            sv = [float(r["saving_pct"]) for r in rows]
            m = sum(sv) / len(sv)
            sd = (sum((x - m) ** 2 for x in sv) / max(1, len(sv) - 1)) ** 0.5
            robust_line = (f"{len(sv)} 县鲁棒：{m:.1f}%±{sd:.1f}"
                           f"（范围 {min(sv):.1f}–{max(sv):.1f}%，全为正）")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_framework(prs)
    slide_algorithm(prs, table_rows)
    slide_application(prs, robust_line)
    prs.save(str(OUT))
    print(f"[OK] {OUT}  ({OUT.stat().st_size/1024:.0f} KB, {len(prs.slides)} slides)")
    print("[table]", table_rows)


if __name__ == "__main__":
    build()
