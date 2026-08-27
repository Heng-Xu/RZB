#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PPT 修订标定：把改动后的文字标红（逐字 diff，只标真正改动的字，未改动字保持原样）。
PowerPoint 无内置修订功能，红字在 Word/WPS/PowerPoint/LibreOffice 一致显示。
不显示删除旧文(避免占位符文本框溢出)。改动规则复用 build_kt_ppt.py。
输出: 徐州公司…研究(定稿-红色标记).pptx
"""
import os, copy, difflib
from pptx import Presentation
from pptx.oxml.ns import qn
from docx.oxml import OxmlElement  # 通用 OOXML 元素构造(命名空间由 tag 决定)
import build_kt_ppt as P  # RULES, ROOT, SRC

OUT = os.path.join(P.ROOT, "研究报告/开题报告/"
      "徐州公司徐州地区分布式新能源高渗透率地区110kV电网容载比弹性指标优化研究(定稿-红色标记).pptx")

def set_red(a_r):
    rPr = a_r.find(qn('a:rPr'))
    if rPr is None:
        rPr = OxmlElement('a:rPr'); a_r.insert(0, rPr)
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        e = rPr.find(qn(tag))
        if e is not None:
            rPr.remove(e)
    sf = OxmlElement('a:solidFill'); clr = OxmlElement('a:srgbClr'); clr.set('val', 'FF0000')
    sf.append(clr)
    ln = rPr.find(qn('a:ln'))
    if ln is not None:
        ln.addnext(sf)
    else:
        rPr.insert(0, sf)

def set_text(a_r, text):
    t = a_r.find(qn('a:t'))
    if t is None:
        t = OxmlElement('a:t'); a_r.append(t)
    t.text = text

def redmark(para, new):
    """按字级 diff 重建 para 的 run：equal 原样、insert/replace 标红新文、delete 略。"""
    ap = para._p
    runs = ap.findall(qn('a:r'))
    if not runs:
        return False
    old = "".join((r.find(qn('a:t')).text or "") for r in runs if r.find(qn('a:t')) is not None)
    tmpl = copy.deepcopy(runs[0])
    endpr = ap.find(qn('a:endParaRPr'))
    for r in runs:
        ap.remove(r)
    def add(seg, red):
        nr = copy.deepcopy(tmpl)
        set_text(nr, seg)
        if red:
            set_red(nr)
        if endpr is not None:
            endpr.addprevious(nr)
        else:
            ap.append(nr)
    sm = difflib.SequenceMatcher(None, old, new, autojunk=False)
    for tag, i1, i2, j1, j2 in sm.get_opcodes():
        if tag == 'equal':
            add(old[i1:i2], False)
        elif tag in ('insert', 'replace'):
            add(new[j1:j2], True)
        # delete: 略过(不显示删除旧文)
    return True

def walk(shapes):
    for sh in shapes:
        try:
            if sh.shape_type == 6 and hasattr(sh, "shapes"):
                yield from walk(sh.shapes); continue
        except Exception:
            pass
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                yield para

def main():
    prs = Presentation(P.SRC)
    done = {i: False for i in range(len(P.RULES))}
    for sl in prs.slides:
        for para in walk(sl.shapes):
            txt = para.text.strip()
            for i, (anc, new) in enumerate(P.RULES):
                if not done[i] and anc in txt:
                    if redmark(para, new):
                        done[i] = True
    try:
        prs.core_properties.author = "XH"
    except Exception:
        pass
    prs.save(OUT)
    print("saved:", OUT, "| slides:", len(prs.slides))
    print("applied:", [i for i in done if done[i]], "| missed:", [i for i in done if not done[i]])

if __name__ == "__main__":
    main()
