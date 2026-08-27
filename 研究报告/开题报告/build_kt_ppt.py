#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""按定稿开题报告口径修正 PPT 文字（结构/页序不变）。
仅做口径对齐: 刚性单值2.0基准 / 双向差异化 / DL/T5729-2023 1.5~2.0 / DL/T2041-2025 反向校核。
就地改段落文字，保留首个 run 字体；不增删形状/页。输出新文件，不覆盖甲方原件。
"""
import os
from pptx import Presentation

ROOT = "/home/roscy/ws_HengXU/徐州地区分布式新能源高渗透率地区110kV电网容载比弹性指标优化研究"
SRC = os.path.join(ROOT, "研究报告/开题报告/徐州公司徐州地区分布式新能源高渗透率地区110kV电网容载比弹性指标优化研究.pptx")
OUT = os.path.join(ROOT, "研究报告/开题报告/徐州公司徐州地区分布式新能源高渗透率地区110kV电网容载比弹性指标优化研究(定稿).pptx")

# (定位子串, 新整段文字)
RULES = [
 ("构提出一片一策",
  "运用层次分析等方法确定新能源增长率、负荷增长率、线路建设成本等指标权重，提出一片一策的弹性容载比规划建议，提升电网规划的科学性、经济性与适应性。"),
 ("严格执行容载比≤2.0和不严格",
  "构建“严格执行容载比≤2.0”与“高渗透片区放宽容载比上限”两类方案的电网建设成本模型（加强10kV联络、配置储能等消纳反送、弥补容量缺口）。"),
 ("场景一：严格执行容载比2.0",
  "场景一（刚性方案）：严格执行容载比≤2.0"),
 ("场景二：弹性容载比策略",
  "场景二（弹性方案）：高渗透片区放宽容载比上限"),
 ("允许一定范围内适度放宽容载比",
  "高渗透率片区适度放宽容载比上限（突破2.0）"),
 ("结合差异化片区策略优化投资",
  "放宽幅度由反向承载力校核与经济性确定"),
 ("依据导则，确定容载比取值弹性区间及边界条件",
  "依据DL/T 5729—2023研究高渗透区上限放宽边界"),
 # 本轮：同步正文精度提升（仅改可换行框，保持简洁）
 ("结合导则中对容载比的最新解读",
  "结合现行导则对容载比口径的最新解读，系统收集统计徐州典型新能源高渗透地区近五年110（35）kV主变负载率、源荷比、电量渗透率及片区联络度等指标，从时间与空间维度开展数据清洗与关联分析，构建以“源荷比—电量渗透率—联络度”三维刻画的“新能源渗透水平—电网运行状态”基础数据库。"),
 ("识别成本驱动因子，采用回归分析等方法",
  "识别主变扩容、线路联络、储能等成本驱动因子，建立参数化电网建设总成本模型"),
 ("量化两种场景的全寿命周期成本，采用年费用法",
  "计入正反向网损、弃光与N-1可靠性，按年费用法量化两类场景全寿命周期成本"),
 ("在安全约束与经济最优间进行多目标求解",
  "在N-1与反向承载力校核约束下进行安全—经济寻优"),
]

def set_para(para, text):
    runs = para.runs
    if not runs:
        r = para.add_run(); r.text = text; return
    runs[0].text = text
    for r in runs[1:]:
        r.text = ""

def walk(shapes):
    for sh in shapes:
        try:
            if sh.shape_type == 6 and hasattr(sh, "shapes"):
                yield from walk(sh.shapes); continue
        except Exception:
            pass
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                yield para

def main():
    prs = Presentation(SRC)
    done = {i: False for i in range(len(RULES))}
    for sl in prs.slides:
        for para in walk(sl.shapes):
            txt = para.text.strip()
            for i, (anc, new) in enumerate(RULES):
                if not done[i] and anc in txt:
                    set_para(para, new); done[i] = True
    try:
        prs.core_properties.author = "XH"
    except Exception:
        pass
    prs.save(OUT)
    print("saved:", OUT, "| slides:", len(prs.slides))
    print("applied:", [i for i in done if done[i]], "| missed:", [i for i in done if not done[i]])

if __name__ == "__main__":
    main()
